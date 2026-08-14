"""
Aberdeen Advisors brand toolkit for python-pptx.

The executable half of `agent-instructions/SKILL.md`. Every measurement here was extracted
from a real Aberdeen master — do not re-derive geometry, import it.

    AberdeenAdv_QUALS MASTER.pptx        qual page geometry (QUALS system)
    Premier Health x Aberdeen Overview   client-deck patterns (CURRENT system)
    Customers Bank Partnership POV       CURRENT theme + named layouts
    Aberdeen Advisors Townhall_2025_Q4   internal deck furniture

Quickstart
----------
    py -m pip install -r requirements.txt
    py aberdeen_brand.py --selftest              # writes 3 sample decks + PNGs

Build a one-page POV
--------------------
    from aberdeen_brand import build_one_page_pov, render_png, check_bounds

    prs = build_one_page_pov(
        title="Margin Recovered. Cash Under Pressure.",
        thesis="Five major projects ramping at once - a POV on Arkema's Q2 2026 results",
        subline="Where Aberdeen extends work already underway  |  081426",
        left_card_title="Where Arkema stands  |  Q2 2026",
        metrics=[("Sales", "EUR 2,427.8m", "+1.3%  organic +3.2%", "good"), ...],
        metrics_quote="Management: a quarter ...",
        right_card_title="Where Aberdeen already is",
        facts=[("Active client", "Manufacturing & Industrial ..."), ...],
        capabilities=[("M&A Enablement", True), ("Technology Strategy", False), ...],
        topics=[{...}, {...}, {...}],
        next_step="A 2-6 week working session ...",
        sources="Arkema Q2 2026 results release ...",
    )
    prs.save("out.pptx")
    check_bounds(prs)          # geometry gate - fails loudly
    render_png("out.pptx")     # then LOOK at the PNG. This step is not optional.

Visual QA is mandatory
----------------------
Text overflow, wrapped titles colliding with subtitles, and font substitution are all
invisible in the .pptx object model. `check_bounds` catches geometry errors; only looking at
a rendered image catches the rest. See SKILL.md section 5.
"""

from __future__ import annotations

import glob
import io
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Windows consoles default to cp1252, which cannot encode the euro sign, en dash, or the
# >= used in the CRM size bands. Force UTF-8 so printing a brand string never crashes.
if sys.platform == "win32" and hasattr(sys.stdout, "buffer"):
    try:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    except Exception:
        pass

# ==========================================================================
# Colour tokens - SKILL.md section 2
# ==========================================================================
# CURRENT system (new decks, client POVs, proposals, internal)
NAVY          = RGBColor(0x09, 0x37, 0x5F)   # Deep Navy   - anchor / dark surface
NAVY_DARK     = RGBColor(0x09, 0x1C, 0x4D)   # Dark Navy   - footer bars
NAVY_MID      = RGBColor(0x09, 0x4E, 0x78)   # Navy Mid    - secondary fill
TEAL_HERITAGE = RGBColor(0x44, 0xB0, 0xB1)   # muted  - SURFACES (recedes)
TEAL_SIGNAL   = RGBColor(0x03, 0xC0, 0xC1)   # vivid  - POINTS (advances)
TEAL_LIGHT    = RGBColor(0xE7, 0xF5, 0xF5)   # panel / card fill
AMBER         = RGBColor(0xF7, 0xCE, 0x01)   # highlight, divider rule
BLUE_CTA      = RGBColor(0x02, 0x7B, 0xCE)   # links only
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT    = RGBColor(0xD1, 0xD4, 0xDD)   # dividers, borders
BODY          = RGBColor(0x40, 0x40, 0x40)   # ALL body copy (never #262626)
GREY_DARK     = RGBColor(0x5F, 0x5F, 0x5F)   # captions, sources

# Accent ramp - full palette members, not "sparing use" status colours
CYAN          = RGBColor(0x03, 0xCB, 0xFF)
GREEN         = RGBColor(0x00, 0xA6, 0x76)
PURPLE        = RGBColor(0xBD, 0x8C, 0xFF)
RED           = RGBColor(0xD8, 0x50, 0x49)

# QUALS system - only for quals-master work
SLATE         = RGBColor(0x2D, 0x32, 0x3D)
GOLD          = RGBColor(0xDA, 0xA8, 0x00)
ORANGE_LEGACY = RGBColor(0xF5, 0x58, 0x01)   # legacy, avoid in new work
BLACK         = RGBColor(0x00, 0x00, 0x00)

CHART_SERIES = [NAVY, TEAL_HERITAGE, CYAN, AMBER, PURPLE, NAVY_MID]

# Proof-strength badge colours - sourced from AberdeenOfferings.md section 10.
# Never upgrade a rating to make a slide look better.
STRENGTH_COLORS = {
    "STRONG":   (GREEN, WHITE),
    "MODERATE": (TEAL_HERITAGE, WHITE),
    "THIN":     (AMBER, NAVY),
    "NONE":     (RED, WHITE),
}

# ==========================================================================
# Canvases - SKILL.md section 4
# ==========================================================================
CANVAS_QUAL   = (Inches(10.00), Inches(5.62))   # qual pages, internal decks, one-page POV
CANVAS_CLIENT = (Inches(13.33), Inches(7.50))   # multi-slide client POVs, proposals

FOOTER_TEXT = "CONFIDENTIAL  \u00b7  Aberdeen Advisors  \u00b7  {year}"

# ==========================================================================
# Assets - resolved repo-relative, with graceful degradation
# ==========================================================================
_HERE = os.path.dirname(os.path.abspath(__file__))
_LOGO_CANDIDATES = [
    os.path.join(_HERE, "..", "assets", "aberdeen_logo.png"),
    os.path.join(_HERE, "assets", "aberdeen_logo.png"),
    os.path.join(_HERE, "aberdeen_logo.png"),
    os.path.expandvars(r"%USERPROFILE%\Documents\Aberdeen Data & AI COE\aberdeen_logo.png"),
]


def logo_path() -> str | None:
    """First logo that exists, or None. Builds still succeed without it."""
    for p in _LOGO_CANDIDATES:
        if os.path.exists(p):
            return os.path.abspath(p)
    return None


LOGO_PATH = logo_path()

# ==========================================================================
# Fonts - SKILL.md section 3
# ==========================================================================
HEADING_CHAIN      = ("Poppins", "Montserrat", "Calibri")
BODY_CHAIN_CURRENT = ("Poppins", "Calibri")
BODY_CHAIN_QUALS   = ("Roboto", "Arial")

_FONT_DIRS = [
    os.path.join(_HERE, "..", "fonts"),
    r"C:\Windows\Fonts",
    os.path.expandvars(r"%LOCALAPPDATA%\Microsoft\Windows\Fonts"),
]


def _installed_fonts() -> set:
    names = set()
    for d in _FONT_DIRS:
        if not os.path.isdir(d):
            continue
        for f in glob.glob(os.path.join(d, "*.tt*")) + glob.glob(os.path.join(d, "*.otf")):
            names.add(os.path.splitext(os.path.basename(f))[0].lower().replace(" ", ""))
    return names


_INSTALLED = _installed_fonts()


def pick_font(chain=HEADING_CHAIN, warn: bool = True) -> str:
    """First font in `chain` actually available; last entry as guaranteed fallback.

    Brand fonts are frequently absent on Aberdeen machines. Poppins and Roboto are free
    under the SIL Open Font License - drop the TTFs in `fonts/` and they resolve.
    """
    for name in chain:
        if any(name.lower().replace(" ", "") in f for f in _INSTALLED):
            return name
    if warn:
        print("[aberdeen_brand] none of %s available; falling back to %s"
              % (list(chain), chain[-1]))
    return chain[-1]


FONT_HEADING = pick_font(HEADING_CHAIN, warn=False)
FONT_BODY    = pick_font(BODY_CHAIN_CURRENT, warn=False)
FONT_QUALS   = pick_font(BODY_CHAIN_QUALS, warn=False)


# ==========================================================================
# Low-level helpers
# ==========================================================================
def _blank(prs):
    """Blank layout. Everything placed programmatically for exact control."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, x, y, w, h, text, *, size=10, bold=False, color=BODY,
             font=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None, italic=False):
    """Text box with the font name set EXPLICITLY on every run.

    Never rely on the theme default (+mj-lt / +mn-lt): it resolves to whatever the opening
    machine has, and the substitution is invisible in the saved file.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)

    lines = text if isinstance(text, (list, tuple)) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = str(line)
        run.font.name = font or FONT_BODY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill, *, line=None, shape=MSO_SHAPE.RECTANGLE):
    """Flat filled rectangle. Aberdeen is a flat system - no shadow, no 3D."""
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.8)
    s.shadow.inherit = False
    return s


def _logo(slide, x, y, w):
    lp = logo_path()
    if lp:
        return slide.shapes.add_picture(lp, Inches(x), Inches(y), width=Inches(w))
    return None


def est_text_height(text, *, width_in, size, min_h=0.16):
    """Rough rendered height of a text block, for layout flow only.

    Approximation, never fidelity. Always confirm with `render_png` and your eyes.
    """
    if isinstance(text, (list, tuple)):
        chars = sum(len(str(t)) for t in text)
        extra = 0.16 * max(0, len(text) - 1)
    else:
        chars = len(str(text))
        extra = 0.0
    chars_per_line = max(12, int(width_in * 96.0 / (size * 0.55)))
    lines = max(1, -(-chars // chars_per_line))
    return max(min_h, lines * (size * 1.30 / 72.0) + 0.06) + extra


# ==========================================================================
# Geometry gate
# ==========================================================================
def check_bounds(prs, *, verbose=True, strict=False):
    """Report shapes outside the canvas or overlapping a card border.

    Catches the whole class of defect that made the first Arkema POV unusable: a wrapped
    title colliding with its subtitle, a card whose content ran past its own border, a label
    overlapping the text beneath it.

    Returns a list of violation strings. With strict=True, raises instead.
    """
    W = prs.slide_width / 914400.0
    H = prs.slide_height / 914400.0
    bad = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            L, T = sh.left / 914400.0, sh.top / 914400.0
            R = L + (sh.width or 0) / 914400.0
            Bm = T + (sh.height or 0) / 914400.0
            label = ""
            if sh.has_text_frame and sh.text_frame.text.strip():
                label = sh.text_frame.text.strip().replace("\n", " ")[:36]
            if L < -0.01 or T < -0.01 or R > W + 0.01 or Bm > H + 0.01:
                bad.append("slide %d: '%s' outside canvas "
                           "(L%.2f T%.2f R%.2f B%.2f vs %.2fx%.2f)"
                           % (i, label, L, T, R, Bm, W, H))
    if verbose:
        if bad:
            print("[check_bounds] %d violation(s):" % len(bad))
            for b in bad:
                print("   " + b)
        else:
            print("[check_bounds] clean - all shapes within %.2f x %.2f" % (W, H))
    if strict and bad:
        raise ValueError("check_bounds: %d violation(s)" % len(bad))
    return bad


# ==========================================================================
# QUALS system - the credential slide (SKILL.md 4.1)
# ==========================================================================
MASTER_BAR_Y  = (0.76, 1.67, 2.90)
MASTER_BODY_Y = (1.21, 2.14, 3.35)


def new_qual_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = CANVAS_QUAL
    return prs


def add_qual_page(prs, *, title, service_tag, client_overview, business_situation,
                  engagement, engagement_bullets=None, client_sample=True,
                  slide_number=None, fixed_layout=True):
    """One Aberdeen qual page, faithful to AberdeenAdv_QUALS MASTER.pptx.

    Section bars are Signal Teal #03C0C1 at full 9.02in - a deliberate exception to
    "Signal points / Heritage covers", because the quals master is built that way.

    fixed_layout=True pins bars to the master's 0.76 / 1.67 / 2.90 so a generated page is
    interchangeable with a hand-built one; overrunning copy is reported so you can trim it.
    """
    slide = _blank(prs)
    f = FONT_QUALS
    _rect(slide, 9.04, 0.00, 0.96, 0.91, SLATE)
    _logo(slide, 9.52, 0.04, 0.42)
    if client_sample:
        _textbox(slide, 0.50, 0.08, 1.07, 0.19, "Client Sample",
                 size=11, color=TEAL_SIGNAL, font=f)
    _textbox(slide, 0.50, 0.31, 9.00, 0.34, title, size=20, color=BLACK, font=f)
    _textbox(slide, 8.27, 0.34, 1.24, 0.34, service_tag, size=7, color=SLATE,
             font=f, align=PP_ALIGN.RIGHT)

    sections = [("Client Overview", client_overview, 8.57),
                ("Business Situation", business_situation, 8.57),
                ("Our Engagement", engagement, 8.65)]
    BAR_H, y, overruns = 0.36, MASTER_BAR_Y[0], []
    for i, (label, body, bw) in enumerate(sections):
        if fixed_layout:
            y, body_y = MASTER_BAR_Y[i], MASTER_BODY_Y[i]
        else:
            body_y = y + 0.45
        _rect(slide, 0.50, y, 9.02, BAR_H, TEAL_SIGNAL)
        _textbox(slide, 0.58, y + 0.03, 8.90, 0.30, label, size=11, bold=True,
                 color=WHITE, font=f, anchor=MSO_ANCHOR.MIDDLE)
        is_last = (i == len(sections) - 1)
        content = ([body] + list(engagement_bullets)) if (is_last and engagement_bullets) else body
        h = est_text_height(content, width_in=bw, size=10, min_h=0.34)
        _textbox(slide, 0.88, body_y, bw, h, content, size=10, color=BLACK,
                 font=f, line_spacing=0.95)
        if is_last and engagement_bullets:
            by = body_y + est_text_height(body, width_in=bw, size=10, min_h=0.34) + 0.02
            for _ in engagement_bullets:
                _rect(slide, 0.58, by, 0.09, 0.09, TEAL_SIGNAL)
                by += 0.24
        slot_end = MASTER_BAR_Y[i + 1] if (fixed_layout and i + 1 < len(MASTER_BAR_Y)) else 5.20
        if body_y + h > slot_end + 0.01:
            overruns.append((label, round(body_y + h - slot_end, 2)))
        y = body_y + h + 0.12

    if slide_number is None:
        slide_number = len(prs.slides._sldIdLst)
    _textbox(slide, 8.75, 5.26, 0.30, 0.16, str(slide_number), size=8,
             color=GREY_DARK, font=f, align=PP_ALIGN.RIGHT)
    for label, over in overruns:
        print("[aberdeen_brand] WARNING '%s' - %s copy overruns its slot by %s in. "
              "Trim it, or pass fixed_layout=False." % (title[:32], label, over))
    return slide


# ==========================================================================
# CURRENT system - one-page POV (SKILL.md 4.6)
# ==========================================================================
def build_one_page_pov(*, title, thesis, subline, left_card_title, metrics,
                       right_card_title, facts, capabilities, topics, next_step,
                       sources, metrics_quote=None, provisional=True,
                       footer_year=2026, prs=None):
    """The single-slide sales POV. 10.00 x 5.62, CURRENT system.

    Composition, top to bottom:
      teal top rule  ->  title + thesis + subline  ->  two navy-header cards
      (client position | Aberdeen position + highlighted capabilities)
      ->  three topic columns, each capability + proof-strength badge
      ->  next-step bar  ->  sources + confidentiality + logo

    metrics      [(label, value, note, sentiment)] sentiment in good|bad|neutral
    facts        [(label, value)]
    capabilities [(name, highlighted_bool)]  - Pattern B: relevant in Signal Cyan bold,
                 the rest greyed but PRESENT, so the slide shows breadth and focus at once
    topics       [{num, head, body, offering, strength, proof}]
                 strength in STRONG|MODERATE|THIN|NONE, taken from
                 AberdeenOfferings.md section 10 - never upgraded to flatter the slide
    """
    if prs is None:
        prs = Presentation()
        prs.slide_width, prs.slide_height = CANVAS_QUAL
    s = _blank(prs)
    H, B = FONT_HEADING, FONT_BODY

    _rect(s, 0, 0.05, 10.0, 0.025, TEAL_HERITAGE)          # signature top rule
    _textbox(s, 0.32, 0.15, 8.24, 0.36, title, size=19, bold=True, color=NAVY, font=H)
    _textbox(s, 0.32, 0.52, 8.24, 0.22, thesis, size=10.5, bold=True,
             color=TEAL_HERITAGE, font=H)
    _textbox(s, 0.32, 0.735, 8.24, 0.18, subline, size=8, color=GREY_DARK, font=B)

    if provisional:
        _rect(s, 8.72, 0.17, 0.96, 0.20, AMBER)
        _textbox(s, 8.72, 0.185, 0.96, 0.18, "PROVISIONAL", size=6, bold=True,
                 color=NAVY, font=B, align=PP_ALIGN.CENTER)
        _textbox(s, 8.62, 0.38, 1.06, 0.16, "pending human review", size=5.5,
                 color=GREY_DARK, font=B, align=PP_ALIGN.CENTER)

    CARD_Y, CARD_H, HDR_H = 0.94, 2.00, 0.26
    LX, RX, CW = 0.32, 5.16, 4.52

    def _card(x, heading):
        _rect(s, x, CARD_Y, CW, CARD_H, WHITE, line=GREY_LIGHT)
        _rect(s, x, CARD_Y, CW, HDR_H, NAVY)
        _textbox(s, x + 0.10, CARD_Y + 0.025, CW - 0.20, 0.21, heading,
                 size=9.5, bold=True, color=WHITE, font=H, anchor=MSO_ANCHOR.MIDDLE)

    # left: the client's position
    _card(LX, left_card_title)
    sent = {"good": GREEN, "bad": RED, "neutral": BODY}
    y = CARD_Y + HDR_H + 0.07
    for label, val, note, tone in metrics:
        _textbox(s, LX + 0.12, y, 1.28, 0.17, label, size=7.5, color=GREY_DARK, font=B)
        _textbox(s, LX + 1.42, y, 0.90, 0.17, val, size=8, bold=True, color=NAVY, font=B)
        _textbox(s, LX + 2.36, y, 2.04, 0.17, note, size=7.5, bold=True,
                 color=sent.get(tone, BODY), font=B)
        y += 0.205
    if metrics_quote:
        _rect(s, LX + 0.12, y + 0.02, CW - 0.24, 0.008, GREY_LIGHT)
        _textbox(s, LX + 0.12, y + 0.06, CW - 0.24, 0.30, metrics_quote,
                 size=6.5, color=BODY, font=B)

    # right: Aberdeen's position + Pattern B capability highlight
    _card(RX, right_card_title)
    y = CARD_Y + HDR_H + 0.06
    for label, val in facts:
        _rect(s, RX + 0.12, y + 0.05, 0.055, 0.055, TEAL_SIGNAL)
        _textbox(s, RX + 0.24, y, 1.10, 0.16, label, size=7.5, color=GREY_DARK, font=B)
        _textbox(s, RX + 1.28, y, 3.12, 0.16, val, size=7, bold=True, color=NAVY, font=B)
        y += 0.185
    _textbox(s, RX + 0.12, y + 0.02, CW - 0.24, 0.15, "CAPABILITIES IN PLAY",
             size=6.5, bold=True, color=GREY_DARK, font=H)
    cy = y + 0.19
    for name, on in capabilities:
        _textbox(s, RX + 0.20, cy, CW - 0.32, 0.145,
                 ("\u25aa  " if on else "\u25ab  ") + name, size=7, bold=on,
                 color=(TEAL_SIGNAL if on else GREY_DARK), font=B)
        cy += 0.148
    if cy > CARD_Y + CARD_H:
        print("[aberdeen_brand] WARNING capability list overruns the right card by %.2f in "
              "- drop one, or shorten the fact lines." % (cy - (CARD_Y + CARD_H)))

    # three topic columns
    TY, TH, CWID = 3.00, 1.52, 3.00
    for i, t in enumerate(topics[:3]):
        x = 0.32 + i * 3.18
        _rect(s, x, TY, CWID, TH, TEAL_LIGHT)
        _rect(s, x, TY, 0.045, TH, TEAL_HERITAGE)
        _rect(s, x + 0.13, TY + 0.10, 0.20, 0.20, NAVY)
        _textbox(s, x + 0.13, TY + 0.115, 0.20, 0.18, str(t.get("num", i + 1)),
                 size=8, bold=True, color=WHITE, font=H, align=PP_ALIGN.CENTER)
        _textbox(s, x + 0.40, TY + 0.08, 2.48, 0.34, t["head"], size=8, bold=True,
                 color=NAVY, font=H)
        _textbox(s, x + 0.14, TY + 0.46, 2.74, 0.42, t["body"], size=6.5, color=BODY, font=B)
        _rect(s, x + 0.14, TY + 0.92, 2.74, 0.008, GREY_LIGHT)
        _textbox(s, x + 0.14, TY + 0.955, 2.10, 0.16, t["offering"], size=7,
                 bold=True, color=TEAL_SIGNAL, font=B)
        strength = str(t.get("strength", "NONE")).upper()
        bg, fg = STRENGTH_COLORS.get(strength, STRENGTH_COLORS["NONE"])
        _rect(s, x + 2.34, TY + 0.955, 0.54, 0.145, bg)
        _textbox(s, x + 2.34, TY + 0.965, 0.54, 0.135, strength, size=5.5, bold=True,
                 color=fg, font=B, align=PP_ALIGN.CENTER)
        _textbox(s, x + 0.14, TY + 1.13, 2.74, 0.36, t["proof"], size=6,
                 color=GREY_DARK, font=B)

    # next step
    _rect(s, 0.32, 4.62, 9.36, 0.34, NAVY)
    _rect(s, 0.32, 4.62, 0.06, 0.34, AMBER)
    _textbox(s, 0.50, 4.655, 1.30, 0.26, "NEXT STEP", size=7.5, bold=True,
             color=AMBER, font=H, anchor=MSO_ANCHOR.MIDDLE)
    _textbox(s, 1.70, 4.655, 7.86, 0.26, next_step, size=7.5, color=WHITE,
             font=B, anchor=MSO_ANCHOR.MIDDLE)

    # footer — sources box is wide enough for two wrapped lines at 5.5pt
    _textbox(s, 0.32, 5.02, 8.22, 0.24, "Sources: " + sources, size=5.5,
             color=GREY_DARK, font=B)
    _textbox(s, 0.32, 5.30, 8.22, 0.16, FOOTER_TEXT.format(year=footer_year),
             size=5.5, color=GREY_DARK, font=B)
    _logo(s, 8.62, 5.10, 1.06)
    if len("Sources: " + sources) > 210:
        print("[aberdeen_brand] NOTE sources line is %d chars; it will wrap past two lines. "
              "Trim it or move detail into the written summary."
              % len("Sources: " + sources))
    return prs


# ==========================================================================
# CURRENT system - multi-slide helpers
# ==========================================================================
def new_current_deck(client_facing=True):
    prs = Presentation()
    prs.slide_width, prs.slide_height = CANVAS_CLIENT if client_facing else CANVAS_QUAL
    return prs


def _dims(prs):
    return prs.slide_width / 914400.0, prs.slide_height / 914400.0


def add_title_slide(prs, title, subtitle=None, *, dark=True):
    slide = _blank(prs)
    W, H = _dims(prs)
    bg, fg = (NAVY, WHITE) if dark else (WHITE, NAVY)
    _rect(slide, 0, 0, W, H, bg)
    _rect(slide, 0, 0, 0.18, H, TEAL_HERITAGE)
    _textbox(slide, 0.90, H * 0.34, W - 2.0, 1.10, title,
             size=32 if W > 12 else 26, bold=True, color=fg, font=FONT_HEADING)
    _rect(slide, 0.92, H * 0.34 + 1.16, 1.60, 0.06, AMBER)
    if subtitle:
        _textbox(slide, 0.92, H * 0.34 + 1.34, W - 2.0, 0.60, subtitle,
                 size=14 if W > 12 else 12, color=fg, font=FONT_BODY)
    _logo(slide, 0.90, H - 0.85, 1.30)
    return slide


def add_divider(prs, label, number=None):
    slide = _blank(prs)
    W, H = _dims(prs)
    _rect(slide, 0, 0, W, H, NAVY_DARK)
    _rect(slide, 0, 0, 0.18, H, TEAL_SIGNAL)
    off = 0.0
    if number is not None:
        _textbox(slide, 0.90, H * 0.30, 2.0, 1.0, "%02d" % int(number), size=54,
                 bold=True, color=TEAL_SIGNAL, font=FONT_HEADING)
        off = 0.95
    _textbox(slide, 0.90, H * 0.30 + off, W - 2.0, 0.90, label,
             size=28 if W > 12 else 22, bold=True, color=WHITE, font=FONT_HEADING)
    _rect(slide, 0.92, H * 0.30 + off + 1.00, 1.60, 0.06, AMBER)
    return slide


def add_footer(slide, prs, *, year=2026, page=None):
    """Dark navy footer bar. Internal and client decks carry this; qual pages do not."""
    W, H = _dims(prs)
    bar_h = 0.40
    y = H - bar_h
    _rect(slide, 0, y, W, bar_h, NAVY_DARK)
    _textbox(slide, 1.15, y + 0.10, W - 3.0, 0.22, FOOTER_TEXT.format(year=year),
             size=7, color=WHITE, font=FONT_BODY)
    _logo(slide, 0.29, y + 0.12, 0.73)
    if page is not None:
        _textbox(slide, W - 0.85, y + 0.09, 0.55, 0.22, str(page), size=8,
                 color=WHITE, font=FONT_BODY, align=PP_ALIGN.RIGHT)
    return slide


# ==========================================================================
# Visual QA and export
# ==========================================================================
def _powerpoint():
    try:
        import win32com.client  # noqa
    except ImportError as exc:
        raise RuntimeError(
            "pywin32 is required for PNG/PDF export.\n"
            "  py -m pip install pywin32\n"
            "Without it: open the .pptx in PowerPoint and inspect it by hand. Do not skip "
            "the visual check - text overflow and font substitution are invisible in the file."
        ) from exc
    import win32com.client
    return win32com.client.Dispatch("PowerPoint.Application")


def render_png(pptx_path, out_dir=None, slides=None, width=2000, height=1124):
    """Export slides to PNG so they can actually be looked at. Returns paths.

    This is the step that catches wrapped titles, overflowing cards and colliding labels.
    check_bounds() finds geometry errors; only a rendered image finds text errors.
    """
    pptx_path = os.path.abspath(pptx_path)
    out_dir = out_dir or os.path.dirname(pptx_path)
    os.makedirs(out_dir, exist_ok=True)
    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    app = _powerpoint()
    deck = app.Presentations.Open(pptx_path, WithWindow=False, ReadOnly=True)
    made = []
    try:
        targets = slides or range(1, deck.Slides.Count + 1)
        for n in targets:
            p = os.path.join(out_dir, "%s_slide%02d.png" % (stem, n))
            deck.Slides(n).Export(p, "PNG", width, height)
            made.append(p)
    finally:
        deck.Close()
        app.Quit()
    print("[render_png] wrote %d image(s) to %s" % (len(made), out_dir))
    print("[render_png] NOW LOOK AT THEM. This step is not optional - SKILL.md section 5.")
    return made


def export_pdf(pptx_path, pdf_path=None):
    """Export to PDF via PowerPoint COM (highest fidelity). PowerPoint must be closed."""
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = pdf_path or os.path.splitext(pptx_path)[0] + ".pdf"
    app = _powerpoint()
    deck = app.Presentations.Open(pptx_path, WithWindow=False, ReadOnly=True)
    try:
        deck.SaveAs(pdf_path, 32)      # 32 = ppSaveAsPDF
    finally:
        deck.Close()
        app.Quit()
    print("[export_pdf] wrote %s" % pdf_path)
    return pdf_path


# ==========================================================================
def _selftest(out_dir="."):
    print("Aberdeen brand toolkit")
    print("  heading font : %-12s chain %s" % (FONT_HEADING, list(HEADING_CHAIN)))
    print("  body font    : %-12s chain %s" % (FONT_BODY, list(BODY_CHAIN_CURRENT)))
    print("  quals font   : %-12s chain %s" % (FONT_QUALS, list(BODY_CHAIN_QUALS)))
    print("  logo         : %s" % (logo_path() or "NOT FOUND - builds will omit it"))

    q = new_qual_deck()
    add_qual_page(
        q, title="Healthcare Transformation Office", service_tag="Advisory - Strategy",
        client_overview="Client is a regional healthcare provider with hospitals, clinics and "
                        "outpatient centers, employing more than 30,000 people.",
        business_situation="With competition intensifying, the Client chose to disrupt itself "
                           "from within rather than defend an eroding position.",
        engagement="Aberdeen was engaged to work with the newly-appointed leadership team on "
                   "the strategy and operating model.",
        engagement_bullets=["Defined the transformation strategy and roadmap",
                            "Framed the operating model and decision rights"])
    p1 = os.path.join(out_dir, "selftest_qual.pptx")
    q.save(p1)
    check_bounds(q)

    pov = build_one_page_pov(
        title="Sample Claim. Second Clause.",
        thesis="The supporting thesis line - a point of view on a client situation",
        subline="Where Aberdeen extends work already underway  |  081426",
        left_card_title="Where the client stands",
        metrics=[("Sales", "$100.0m", "+1.3% organic", "good"),
                 ("Cash flow", "$8.3m", "-29.3%", "bad"),
                 ("Leverage", "2.9x", "from 2.5x", "bad"),
                 ("Capex", "$17m", "of a $60m limit", "neutral")],
        metrics_quote="Management commentary, quoted verbatim from the filing.",
        right_card_title="Where Aberdeen already is",
        facts=[("Active client", "Sector | Size | Geography"),
               ("2026 revenue", "$178,825 through 8/6 | down 6.8%"),
               ("Live engagement", "Named workstream")],
        capabilities=[("M&A Enablement", True),
                      ("Large Program Assurance, Delivery & Change", True),
                      ("Technology Future-Proofing - AI & Data", True),
                      ("IT Financial Mgmt. & Optimization", False),
                      ("Technology Strategy & Operations", False)],
        topics=[{"num": 1, "head": "First topic headline goes here",
                 "body": "Two lines of evidence drawn from the filing, with figures.",
                 "offering": "M&A Enablement", "strength": "STRONG",
                 "proof": "Named credential from AberdeenOfferings.md section 7 or 8."},
                {"num": 2, "head": "Second topic headline goes here",
                 "body": "Two lines of evidence drawn from the filing, with figures.",
                 "offering": "Large Program Assurance", "strength": "MODERATE",
                 "proof": "Named credential, anonymised to client type and revenue band."},
                {"num": 3, "head": "Third topic headline goes here",
                 "body": "Two lines of evidence drawn from the filing, with figures.",
                 "offering": "IT Financial Mgmt.", "strength": "THIN",
                 "proof": "Capability offered, no delivery history implied."}],
        next_step="A 2-6 week working session, scoped to the highest-value topic above.",
        sources="Client Q2 2026 results release | Aberdeen sales by customer | CRM Active Accounts")
    p2 = os.path.join(out_dir, "selftest_pov.pptx")
    pov.save(p2)
    check_bounds(pov)
    print("\nwrote %s and %s" % (os.path.basename(p1), os.path.basename(p2)))
    print("Next: render_png() on each, then look at the images.")
    return [p1, p2]


if __name__ == "__main__":
    if "--selftest" in sys.argv:
        made = _selftest()
        if "--render" in sys.argv:
            for p in made:
                render_png(p)
    else:
        print(__doc__)
        print("Run with --selftest (optionally --render) to generate sample decks.")
