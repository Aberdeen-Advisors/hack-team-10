"""Build the 4-slide Pursuit Intelligence summary deck.

Brand tokens taken from the app's own styles.css (:root), which is styled per
About Aberdeen/Aberdeen Slide Template_July 2025.pptx — so the deck and the app
match. Note the app's signal accent is #03CBFF (accent2), not the #03C0C1 used
by the legacy quals master.

Screenshots are real captures of the running app, cropped to remove the nav bar.
"""
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SHOTS = sys.argv[1]
OUT = sys.argv[2]
LOGO = sys.argv[3] if len(sys.argv) > 3 else None

# --- tokens, verbatim from app/styles.css :root ---------------------------
NAVY       = RGBColor(0x09, 0x37, 0x5F)
NAVY_DARK  = RGBColor(0x04, 0x1B, 0x2F)
TEAL_HER   = RGBColor(0x44, 0xB0, 0xB1)
TEAL_SIG   = RGBColor(0x03, 0xCB, 0xFF)
TEAL_LIGHT = RGBColor(0xE7, 0xF5, 0xF5)
AMBER      = RGBColor(0xF7, 0xCE, 0x01)
BLUE_CTA   = RGBColor(0x00, 0x72, 0xAD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LIGHT = RGBColor(0xD1, 0xD4, 0xDD)
BODY       = RGBColor(0x40, 0x40, 0x40)
GREY_DARK  = RGBColor(0x5F, 0x5F, 0x5F)
GREEN      = RGBColor(0x00, 0xA6, 0x76)
PURPLE     = RGBColor(0xBD, 0x8C, 0xFF)
BG         = RGBColor(0xF7, 0xF7, 0xFC)
F = "Poppins"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(10.00), Inches(5.62)


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    o = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line is None:
        o.line.fill.background()
    else:
        o.line.color.rgb = line; o.line.width = Pt(0.75)
    o.shadow.inherit = False
    return o


def tb(s, x, y, w, h, text, size=10, bold=False, color=BODY, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, italic=False, spacing=None):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(text if isinstance(text, (list, tuple)) else [text]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        r = p.add_run(); r.text = str(line)
        r.font.name = F; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
    return box


def chrome(s, title, kicker=None, page=None):
    """Teal top rule, title, optional kicker, footer with logo and page number."""
    rect(s, 0, 0, 10.0, 0.055, TEAL_HER)
    tb(s, 0.34, 0.16, 8.5, 0.40, title, size=20, bold=True, color=NAVY)
    if kicker:
        tb(s, 0.34, 0.575, 8.9, 0.20, kicker, size=8.5, color=GREY_DARK)
    if LOGO and os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(8.70), Inches(5.16), width=Inches(0.98))
    tb(s, 0.34, 5.24, 5.0, 0.16, "CONFIDENTIAL  ·  Aberdeen Advisors  ·  2026",
       size=5.5, color=GREY_DARK)
    if page:
        tb(s, 8.10, 5.24, 0.45, 0.16, str(page), size=7, color=GREY_DARK, align=PP_ALIGN.RIGHT)


def crop(name, box, out_name):
    src = os.path.join(SHOTS, name)
    im = Image.open(src)
    im.crop(box).save(os.path.join(SHOTS, out_name))
    return os.path.join(SHOTS, out_name)


def fit_pic(s, path, x, y, max_w, max_h):
    """Place a picture scaled to fit inside the box, preserving aspect ratio.

    Sizing by width alone is how the first build of this deck overflowed: a
    2410x1310 crop at 9.2in wide renders 5.0in tall, off the bottom of a 5.62in
    slide. Always constrain both axes.
    """
    im = Image.open(path)
    ratio = im.width / im.height
    w, h = max_w, max_w / ratio
    if h > max_h:
        h, w = max_h, max_h * ratio
    return s.shapes.add_picture(path, Inches(x + (max_w - w) / 2),
                               Inches(y + (max_h - h) / 2),
                               width=Inches(w), height=Inches(h))


def kpi(s, x, y, w, num, label, accent=TEAL_SIG):
    rect(s, x, y, w, 0.62, NAVY)
    rect(s, x, y, 0.045, 0.62, accent)
    tb(s, x + 0.14, y + 0.05, w - 0.22, 0.30, num, size=17, bold=True, color=accent)
    tb(s, x + 0.14, y + 0.36, w - 0.22, 0.22, label, size=7, color=WHITE)


# =========================================================================
# SLIDE 1 — Overview
# =========================================================================
s1 = blank()
rect(s1, 0, 0, 10.0, 5.62, BG)
chrome(s1, "Pursuit Intelligence", None, 1)
tb(s1, 0.34, 0.56, 5.60, 0.22,
   "Public filings plus Aberdeen's own credentials, as a ranked pursuit list",
   size=9, bold=True, color=TEAL_HER)
tb(s1, 0.34, 0.83, 5.55, 0.62,
   "A working web app, built entirely from what is already in the repo: three agent-instruction "
   "files and 13 public company documents. It scores every account, shows why, and never invents "
   "a fact it cannot cite.",
   size=8, color=BODY, spacing=1.06)

for i, (n, l, c) in enumerate([("13", "Accounts assessed", TEAL_SIG),
                               ("59", "Evidenced buying signals", TEAL_SIG),
                               ("33", "Addressable initiatives", AMBER),
                               ("0", "External API calls", GREEN)]):
    kpi(s1, 0.34 + i * 1.40, 1.52, 1.28, n, l, c)

# slim architecture ribbon
tb(s1, 0.34, 2.32, 5.55, 0.16, "HOW IT WORKS", size=6.5, bold=True, color=GREY_DARK)
steps = ["13 public\ndocuments", "59 signals\nextracted", "Scored\n40/30/30",
         "Ranked\npipeline", "Pursuit brief\n+ outreach"]
for i, st in enumerate(steps):
    x = 0.34 + i * 1.13
    rect(s1, x, 2.52, 1.00, 0.46, WHITE, line=GREY_LIGHT)
    tb(s1, x + 0.04, 2.575, 0.92, 0.36, st.replace("\n", "  "), size=6, bold=True,
       color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        tb(s1, x + 1.00, 2.60, 0.13, 0.30, "›", size=11, bold=True, color=TEAL_HER,
           align=PP_ALIGN.CENTER)

tb(s1, 0.34, 3.14, 5.55, 0.20, "WHAT MAKES IT DEFENSIBLE", size=6.5, bold=True, color=GREY_DARK)
for i, (h, d) in enumerate([
        ("Every claim cites a page", "Signals carry the document and page reference they came from."),
        ("Proof ratings are never upgraded", "STRONG / MODERATE / THIN come from AberdeenOfferings.md §10."),
        ("Gaps are reported, not filled", "Missing CRM fields show as gaps; one stale source is excluded."),
        ("Nothing leaves the machine", "Level 3 compliant: zero network calls anywhere in the app.")]):
    y = 3.34 + i * 0.42
    rect(s1, 0.36, y + 0.055, 0.075, 0.075, TEAL_SIG)
    tb(s1, 0.52, y, 5.36, 0.18, h, size=7.5, bold=True, color=NAVY)
    tb(s1, 0.52, y + 0.165, 5.36, 0.20, d, size=6.5, color=GREY_DARK)

# full-page capture so all 13 ranked rows are visible, cropped to the table
shot = crop("tab_2_pipeline.png", (250, 600, 2660, 2620), "c_pipeline.png")
rect(s1, 5.98, 0.84, 3.70, 4.02, WHITE, line=GREY_LIGHT)
fit_pic(s1, shot, 6.04, 0.90, 3.58, 3.90)
tb(s1, 5.98, 4.90, 3.70, 0.16, "All 13 accounts ranked — click any row for the full brief",
   size=6, italic=True, color=GREY_DARK, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2 — User benefit
# =========================================================================
s2 = blank()
rect(s2, 0, 0, 10.0, 5.62, BG)
chrome(s2, "One brief, three sources, one recommendation", None, 2)
tb(s2, 0.34, 0.56, 9.3, 0.22,
   "What a partner previously assembled by hand over hours — now on one screen, with every "
   "figure traceable",
   size=9, bold=True, color=TEAL_HER)

cols = [("WHAT WE ALREADY KNOW", TEAL_HER,
         ["Relationship position and play", "Live engagement, if any",
          "Revenue trajectory and bucket", "Blank CRM fields shown as gaps"]),
        ("WHAT THE RESEARCH ADDS", TEAL_SIG,
         ["Forward investments from filings", "Committed vs announced vs exploratory",
          "Decision-makers and topic to raise", "Page-level citation on every claim"]),
        ("WHAT WE CAN SELL", AMBER,
         ["Offerings the signals trigger", "Credentials, same-sector first",
          "Honest proof strength badge", "Draft outreach, partner-editable"])]
for i, (head, accent, items) in enumerate(cols):
    x = 0.34 + i * 3.12
    rect(s2, x, 0.86, 2.94, 1.06, WHITE, line=GREY_LIGHT)
    rect(s2, x, 0.86, 2.94, 0.045, accent)
    tb(s2, x + 0.10, 0.925, 2.74, 0.16, head, size=6.5, bold=True, color=NAVY)
    for j, it in enumerate(items):
        rect(s2, x + 0.12, 1.155 + j * 0.185, 0.055, 0.055, accent)
        tb(s2, x + 0.24, 1.105 + j * 0.185, 2.62, 0.18, it, size=6, color=BODY)

# wide, short crop of the brief so it fits the band without overflowing
# bottom edge lands just below the Power BI credential block, so the crop ends on
# a natural boundary rather than slicing a heading in half
shot2 = crop("brief_fold.png", (250, 480, 2660, 1205), "c_brief.png")
rect(s2, 0.34, 2.00, 9.32, 2.86, WHITE, line=GREY_LIGHT)
fit_pic(s2, shot2, 0.40, 2.06, 9.20, 2.74)
rect(s2, 0.34, 4.90, 9.32, 0.26, NAVY)
tb(s2, 0.44, 4.925, 9.1, 0.22,
   "Arkema brief — filings research on the left, Aberdeen capability and proof on the right, "
   "with the live Clear Lake engagement flagged as in progress, never as a delivered outcome",
   size=6.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# SLIDE 3 — Architecture
# =========================================================================
s3 = blank()
rect(s3, 0, 0, 10.0, 5.62, BG)
chrome(s3, "Architecture", None, 3)
tb(s3, 0.34, 0.56, 9.3, 0.22,
   "Static, self-contained, zero egress — the whole app is one HTML file with the data compiled in",
   size=9, bold=True, color=TEAL_HER)

def layer(y, h, title, sub, fill, tcol=WHITE, accent=None):
    rect(s3, 0.34, y, 9.32, h, fill)
    if accent: rect(s3, 0.34, y, 0.05, h, accent)
    tb(s3, 0.50, y + 0.05, 2.30, 0.20, title, size=8, bold=True, color=tcol)
    tb(s3, 0.50, y + 0.235, 2.30, 0.18, sub, size=6, color=tcol)

# inputs
layer(0.86, 0.50, "1 · SOURCES", "committed to the repo", NAVY, WHITE, TEAL_HER)
for i, (t, d) in enumerate([("client-data/", "13 public documents — 9 press releases, 3 filings, 1 digest"),
                            ("agent-instructions/", "AGENT_WORKFLOW · SKILL · AberdeenOfferings"),
                            ("About Aberdeen/", "quals master + enhanced credential library")]):
    x = 3.00 + i * 2.24
    rect(s3, x, 0.92, 2.14, 0.38, WHITE, line=GREY_LIGHT)
    tb(s3, x + 0.08, 0.955, 1.98, 0.16, t, size=6.5, bold=True, color=NAVY)
    tb(s3, x + 0.08, 1.105, 1.98, 0.16, d, size=5.5, color=GREY_DARK)

tb(s3, 4.9, 1.40, 0.4, 0.18, "▼", size=8, color=TEAL_HER, align=PP_ALIGN.CENTER)

# data layer
layer(1.60, 0.72, "2 · DATA LAYER", "hand-extracted, compiled in", NAVY_DARK, WHITE, TEAL_SIG)
for i, (t, d) in enumerate([("data-signals.js", "59 signals · evidence, page ref, provenance, category, triggered offering"),
                            ("data-investments.js", "33 addressable initiatives · committed / announced / exploratory"),
                            ("data-contacts.js", "decision-makers with the topic to raise"),
                            ("data-aberdeen.js", "offerings, credentials, §10 proof ratings, account buckets")]):
    x = 3.00 + (i % 2) * 3.36
    y = 1.66 + (i // 2) * 0.31
    rect(s3, x, y, 3.26, 0.27, WHITE, line=GREY_LIGHT)
    tb(s3, x + 0.07, y + 0.015, 1.10, 0.14, t, size=6, bold=True, color=NAVY)
    tb(s3, x + 1.18, y + 0.015, 2.02, 0.24, d, size=5, color=GREY_DARK)

tb(s3, 4.9, 2.40, 0.4, 0.18, "▼", size=8, color=TEAL_HER, align=PP_ALIGN.CENTER)

# logic
layer(2.60, 0.62, "3 · SCORING", "transparent, arguable", NAVY, WHITE, AMBER)
for i, (t, d, mx) in enumerate([("Signal evidence", "count, weighted for tech / AI / M&A / restructuring", "40"),
                                ("Relationship position", "the workflow bucket — live engagement scores highest", "30"),
                                ("Proof-strength fit", "mean §10 rating across triggered offerings", "30")]):
    x = 3.00 + i * 2.24
    rect(s3, x, 2.66, 2.14, 0.50, TEAL_LIGHT, line=GREY_LIGHT)
    tb(s3, x + 0.08, 2.685, 1.55, 0.16, t, size=6.5, bold=True, color=NAVY)
    tb(s3, x + 1.66, 2.675, 0.42, 0.18, mx, size=9, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    tb(s3, x + 0.08, 2.855, 1.98, 0.28, d, size=5, color=GREY_DARK)

tb(s3, 4.9, 3.26, 0.4, 0.18, "▼", size=8, color=TEAL_HER, align=PP_ALIGN.CENTER)

# outputs
layer(3.46, 0.50, "4 · SURFACES", "three views, one file", NAVY, WHITE, GREEN)
for i, (t, d) in enumerate([("Radar", "where capital is going, filtered to Aberdeen-addressable"),
                            ("Pipeline", "13 accounts ranked, tiered Pursue / Develop / Qualify"),
                            ("Sources", "every document, provenance tier and guardrail")]):
    x = 3.00 + i * 2.24
    rect(s3, x, 3.52, 2.14, 0.38, WHITE, line=GREY_LIGHT)
    tb(s3, x + 0.08, 3.555, 1.98, 0.16, t, size=6.5, bold=True, color=NAVY)
    tb(s3, x + 0.08, 3.705, 1.98, 0.16, d, size=5.5, color=GREY_DARK)

# guardrail band
rect(s3, 0.34, 4.10, 9.32, 0.60, TEAL_LIGHT, line=TEAL_HER)
tb(s3, 0.50, 4.16, 2.20, 0.18, "GUARDRAILS — ENFORCED, NOT ASPIRATIONAL", size=6.5, bold=True, color=NAVY)
for i, g in enumerate(["Zero network calls · Level 3 stays local",
                       "Page-level citation on every figure",
                       "Proof ratings never upgraded",
                       "Stale source excluded, not mined",
                       "Unconfirmed name matches flagged",
                       "Blank CRM fields reported as gaps"]):
    x = 0.50 + (i % 3) * 3.06
    y = 4.36 + (i // 3) * 0.16
    rect(s3, x, y + 0.045, 0.05, 0.05, TEAL_HER)
    tb(s3, x + 0.12, y - 0.005, 2.90, 0.15, g, size=5.5, color=BODY)

tb(s3, 0.34, 4.78, 9.32, 0.18,
   "Deliberate trade-off: the data layer is compiled by hand, so nothing drifts mid-demo and there is "
   "no credential to leak — but it does not refresh itself. That is the first thing to change at scale.",
   size=6, italic=True, color=GREY_DARK)

# =========================================================================
# SLIDE 4 — Path to scale
# =========================================================================
s4 = blank()
rect(s4, 0, 0, 10.0, 5.62, BG)
chrome(s4, "Path to scale", None, 4)
tb(s4, 0.34, 0.56, 9.3, 0.22,
   "From 13 hand-compiled accounts to the firm's whole book, without giving up the evidence discipline",
   size=9, bold=True, color=TEAL_HER)

phases = [
    ("NOW", "Proof of concept", TEAL_HER,
     ["13 accounts, 59 signals, hand-extracted",
      "Three views, one self-contained file",
      "Runs with no install and no credential",
      "Brand-correct, Poppins embedded"]),
    ("30 DAYS", "Connect the real book", TEAL_SIG,
     ["Read the 30-account CRM and revenue file",
      "Join credentials.py to the quals library",
      "Generate the one-page POV from the brief",
      "Fill the 48 missing quantified outcomes"]),
    ("60 DAYS", "Automate the intake", AMBER,
     ["Scheduled filing and news collection",
      "Signal recency decay and expiry",
      "Sort, filter and search the pipeline",
      "Partner override with an audit trail"]),
    ("90 DAYS", "Make it the BD operating system", GREEN,
     ["Warm-path graph from engagement history",
      "Outreach tracked through to reply",
      "Win/loss written back as new credentials",
      "Feeds HorizonView, not competes with it"]),
]
for i, (when, title, accent, items) in enumerate(phases):
    x = 0.34 + i * 2.36
    rect(s4, x, 0.90, 2.22, 2.56, WHITE, line=GREY_LIGHT)
    rect(s4, x, 0.90, 2.22, 0.30, NAVY)
    tb(s4, x + 0.10, 0.925, 1.0, 0.22, when, size=7, bold=True, color=accent)
    tb(s4, x + 0.10, 1.27, 2.02, 0.34, title, size=8.5, bold=True, color=NAVY)
    for j, it in enumerate(items):
        rect(s4, x + 0.12, 1.73 + j * 0.40, 0.055, 0.055, accent)
        tb(s4, x + 0.25, 1.675 + j * 0.40, 1.86, 0.38, it, size=6, color=BODY, spacing=1.02)
    if i < 3:
        tb(s4, x + 2.22, 1.95, 0.14, 0.24, "›", size=12, bold=True, color=TEAL_HER,
           align=PP_ALIGN.CENTER)

# the three real blockers
rect(s4, 0.34, 3.60, 9.32, 1.06, WHITE, line=GREY_LIGHT)
rect(s4, 0.34, 3.60, 0.05, 1.06, AMBER)
tb(s4, 0.50, 3.66, 4.0, 0.18, "WHAT ACTUALLY GATES THIS", size=6.5, bold=True, color=NAVY)
for i, (h, d) in enumerate([
        ("48 blank outcomes", "Four of six capabilities have zero quantified proof. One number per qual converts them to sellable."),
        ("One source of truth", "The data layer duplicates the instruction files. Generate it from them instead."),
        ("Local-path dependency", "The workflow assumes one machine. Repo-relative paths let anyone run it.")]):
    y = 3.88 + i * 0.25
    rect(s4, 0.52, y + 0.05, 0.05, 0.05, AMBER)
    tb(s4, 0.66, y - 0.01, 1.62, 0.16, h, size=6.5, bold=True, color=NAVY)
    tb(s4, 2.34, y - 0.01, 7.20, 0.16, d, size=6, color=GREY_DARK)

rect(s4, 0.34, 4.76, 9.32, 0.30, NAVY)
tb(s4, 0.44, 4.79, 9.1, 0.24,
   "The discipline is the asset: a tool a partner can argue with beats one they have to trust.",
   size=7.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("wrote %s  (%d slides)" % (os.path.basename(OUT), len(prs.slides._sldIdLst)))
